from pathlib import Path

from pptx import Presentation as PptxPresentation

from src.models import Presentation, Slide, SlideType, TableData, Metadata
from src.renderer import PptxRenderer
from src.theme import load_theme


def _make_presentation(*slides: Slide) -> Presentation:
    return Presentation(
        title="Test",
        slides=list(slides),
        metadata=Metadata(page_count=len(slides)),
    )


def test_render_creates_valid_pptx(tmp_path):
    slide = Slide(number=1, title="Hello", type=SlideType.TITLE)
    pres = _make_presentation(slide)
    renderer = PptxRenderer()
    out = tmp_path / "test.pptx"
    renderer.render(pres, out)
    assert out.exists()
    pptx = PptxPresentation(str(out))
    assert len(pptx.slides) >= 1


def test_render_title_slide(tmp_path):
    slide = Slide(number=1, title="My Title", type=SlideType.TITLE, body=["v1.0"])
    pres = _make_presentation(slide)
    out = tmp_path / "title.pptx"
    PptxRenderer().render(pres, out)
    pptx = PptxPresentation(str(out))
    assert len(pptx.slides) >= 1


def test_render_concept_with_bullets(tmp_path):
    slide = Slide(
        number=1, title="Concept", type=SlideType.CONCEPT,
        body=["● Point 1", "  ○ Sub point", "● Point 2"],
    )
    pres = _make_presentation(slide)
    out = tmp_path / "concept.pptx"
    PptxRenderer().render(pres, out)
    pptx = PptxPresentation(str(out))
    assert len(pptx.slides) >= 1


def test_render_table_in_body(tmp_path):
    table = TableData(headers=["Name", "Value"], rows=[["A", "1"], ["B", "2"]])
    slide = Slide(number=1, title="Data", type=SlideType.CONCEPT, body=["Intro text", table])
    pres = _make_presentation(slide)
    out = tmp_path / "table.pptx"
    PptxRenderer().render(pres, out)
    pptx = PptxPresentation(str(out))
    assert len(pptx.slides) >= 1


def test_render_speaker_notes(tmp_path):
    slide = Slide(number=1, title="Notes Test", type=SlideType.CONCEPT, speaker_note="This is a speaker note")
    pres = _make_presentation(slide)
    out = tmp_path / "notes.pptx"
    PptxRenderer().render(pres, out)
    pptx = PptxPresentation(str(out))
    notes = pptx.slides[0].notes_slide.notes_text_frame.text
    assert "speaker note" in notes


def test_section_divider_inserted(tmp_path):
    s1 = Slide(number=1, title="A", type=SlideType.CONCEPT, section="Sec1")
    s2 = Slide(number=2, title="B", type=SlideType.CONCEPT, section="Sec2")
    pres = _make_presentation(s1, s2)
    out = tmp_path / "sections.pptx"
    PptxRenderer().render(pres, out)
    pptx = PptxPresentation(str(out))
    assert len(pptx.slides) == 3  # s1 + divider + s2


def test_image_placeholder_slide(tmp_path):
    slide = Slide(number=1, title="Chart", type=SlideType.IMAGE_REFERENCE, image_placeholder=True, body=["Session/PV 비교 차트"])
    pres = _make_presentation(slide)
    out = tmp_path / "image.pptx"
    PptxRenderer().render(pres, out)
    pptx = PptxPresentation(str(out))
    assert len(pptx.slides) >= 1
