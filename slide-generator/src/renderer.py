from __future__ import annotations

import logging
import math
import re
from pathlib import Path

from pptx import Presentation as PptxPresentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from src.models import Presentation, Slide, SlideType, TableData
from src.theme import Theme, load_theme

logger = logging.getLogger(__name__)

# Slide types that use generic content rendering
_CONTENT_TYPES = {
    SlideType.CONCEPT,
    SlideType.OVERVIEW,
    SlideType.FRAMEWORK,
    SlideType.DATA_REPORT,
    SlideType.APPENDIX,
}

# Height estimation constants (inches)
BULLET_HEIGHT = 0.32
TABLE_ROW_HEIGHT = 0.28
TABLE_PADDING = 0.15
SLIDE_BOTTOM_MARGIN = 0.6


class PptxRenderer:
    """Renders a Presentation model to a .pptx file with Memphis Corporate style."""

    def __init__(self, theme: Theme | None = None) -> None:
        self.theme = theme or load_theme()

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, presentation: Presentation, output_path: Path) -> Path:
        pptx = PptxPresentation()
        pptx.slide_width = Inches(self.theme.meta.slide_width)
        pptx.slide_height = Inches(self.theme.meta.slide_height)

        blank_layout = pptx.slide_layouts[6]
        prev_section: str = ""

        for slide_model in presentation.slides:
            # Auto-insert section divider when section changes
            if (
                slide_model.section
                and prev_section
                and slide_model.section != prev_section
            ):
                self._add_section_divider(pptx, blank_layout, slide_model.section)

            # Auto-split: split body into multiple slides if overflowing
            sub_slides = self._split_slide_if_needed(slide_model)

            for idx, sub in enumerate(sub_slides):
                pptx_slide = pptx.slides.add_slide(blank_layout)
                self._apply_background(pptx_slide)
                self._dispatch(pptx_slide, sub)

                # Speaker notes only on the first sub-slide
                if idx == 0 and slide_model.speaker_note:
                    notes_slide = pptx_slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_model.speaker_note

            prev_section = slide_model.section

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        pptx.save(str(output_path))
        return output_path

    # ------------------------------------------------------------------
    # Auto-split logic
    # ------------------------------------------------------------------

    def _estimate_item_height(self, item) -> float:
        """Estimate the rendered height of a single body item in inches."""
        if isinstance(item, TableData):
            return (len(item.rows) + 1) * TABLE_ROW_HEIGHT + TABLE_PADDING
        text = str(item)
        # Rough estimate: long lines wrap, so add extra height
        line_count = max(1, math.ceil(len(text) / 70))
        return BULLET_HEIGHT * line_count

    def _max_body_height(self) -> float:
        """Available height for body content on a single slide."""
        return self._slide_h - self.theme.layout.body_top - SLIDE_BOTTOM_MARGIN

    def _split_slide_if_needed(self, slide: Slide) -> list[Slide]:
        """Split a slide into multiple if body overflows available space."""
        if not slide.body:
            return [slide]
        if slide.type in (SlideType.TITLE, SlideType.IMAGE_REFERENCE, SlideType.SUMMARY):
            return [slide]

        max_h = self._max_body_height()
        pages: list[list] = []
        current_page: list = []
        current_h = 0.0

        for item in slide.body:
            item_h = self._estimate_item_height(item)
            if current_page and (current_h + item_h) > max_h:
                pages.append(current_page)
                current_page = []
                current_h = 0.0
            current_page.append(item)
            current_h += item_h

        if current_page:
            pages.append(current_page)

        if len(pages) <= 1:
            return [slide]

        # Create sub-slides
        result: list[Slide] = []
        for i, page_body in enumerate(pages):
            suffix = f" ({i + 1}/{len(pages)})" if len(pages) > 1 else ""
            result.append(slide.model_copy(update={
                "title": slide.title + suffix,
                "body": page_body,
            }))
        return result

    # ------------------------------------------------------------------
    # Dispatcher
    # ------------------------------------------------------------------

    def _dispatch(self, pptx_slide, slide: Slide) -> None:
        if slide.type == SlideType.TITLE:
            self._render_title_slide(pptx_slide, slide)
        elif slide.type == SlideType.SUMMARY:
            self._render_summary_slide(pptx_slide, slide)
        elif slide.type == SlideType.EXERCISE:
            self._render_exercise_slide(pptx_slide, slide)
        elif slide.type == SlideType.EXAMPLE:
            self._render_example_slide(pptx_slide, slide)
        elif slide.type == SlideType.CRM_EXPERIMENT:
            self._render_two_column_slide(pptx_slide, slide)
        elif slide.type == SlideType.IMAGE_REFERENCE:
            self._render_image_placeholder(pptx_slide, slide)
        else:
            self._render_content_slide(pptx_slide, slide)

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @property
    def _slide_w(self) -> float:
        return self.theme.meta.slide_width

    @property
    def _slide_h(self) -> float:
        return self.theme.meta.slide_height

    @property
    def _margin(self):
        return self.theme.layout.margin

    @property
    def _usable_width(self) -> float:
        return self._slide_w - self._margin.left - self._margin.right

    def _add_textbox(self, slide, left, top, width, height):
        return slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height),
        )

    def _set_font(self, run, font_cfg, color_rgb: tuple[int, int, int] | None = None):
        run.font.name = font_cfg.name
        run.font.size = Pt(font_cfg.size)
        run.font.bold = font_cfg.bold
        if color_rgb:
            run.font.color.rgb = RGBColor(*color_rgb)

    def _apply_background(self, slide):
        """Apply light background color to slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        r, g, b = self.theme.colors.hex_to_rgb("background")
        fill.fore_color.rgb = RGBColor(r, g, b)

    # ------------------------------------------------------------------
    # Memphis decorative elements
    # ------------------------------------------------------------------

    def _add_accent_bar(self, slide, top: float, width: float = 0.8, color_name: str = "accent"):
        """Add a small colored accent bar (Memphis style)."""
        r, g, b = self.theme.colors.hex_to_rgb(color_name)
        shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            Inches(self._margin.left), Inches(top),
            Inches(width), Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()
        return shape

    def _add_side_decoration(self, slide, color_name: str = "accent_yellow"):
        """Add a subtle side decoration strip (Memphis pattern element)."""
        r, g, b = self.theme.colors.hex_to_rgb(color_name)
        shape = slide.shapes.add_shape(
            1,  # RECTANGLE
            Inches(self._slide_w - 0.15), Inches(0),
            Inches(0.15), Inches(self._slide_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()

    def _add_corner_dot(self, slide, left: float, top: float, size: float = 0.18, color_name: str = "accent"):
        """Add a small decorative circle (Memphis dot)."""
        r, g, b = self.theme.colors.hex_to_rgb(color_name)
        shape = slide.shapes.add_shape(
            9,  # OVAL
            Inches(left), Inches(top),
            Inches(size), Inches(size),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()

    def _add_page_number(self, slide, slide_model: Slide):
        """Add page number at bottom right."""
        t = self.theme
        txbox = self._add_textbox(
            slide,
            self._slide_w - self._margin.right - 0.8,
            self._slide_h - 0.4,
            0.8, 0.3,
        )
        tf = txbox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = str(slide_model.number)
        run.font.name = t.fonts.note.name
        run.font.size = Pt(t.fonts.note.size)
        r, g, b = t.colors.hex_to_rgb("secondary")
        run.font.color.rgb = RGBColor(r, g, b)

    # ------------------------------------------------------------------
    # Title rendering with accent bar
    # ------------------------------------------------------------------

    def _add_title(self, slide, title: str, *, top: float | None = None,
                   align=PP_ALIGN.LEFT, slide_model: Slide | None = None):
        t = self.theme
        top = top if top is not None else t.layout.title_top

        # Accent bar above title
        self._add_accent_bar(slide, top - 0.12, width=0.6)

        txbox = self._add_textbox(
            slide, self._margin.left, top, self._usable_width, 0.7,
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = title
        self._set_font(run, t.fonts.title, t.colors.hex_to_rgb("primary"))

        # Section subtitle if available
        if slide_model and slide_model.section:
            sp = tf.add_paragraph()
            sp.alignment = align
            sr = sp.add_run()
            sr.text = slide_model.section
            sr.font.name = t.fonts.subtitle.name
            sr.font.size = Pt(t.fonts.subtitle.size - 2)
            r, g, b = t.colors.hex_to_rgb("secondary")
            sr.font.color.rgb = RGBColor(r, g, b)

        return txbox

    # ------------------------------------------------------------------
    # Body rendering
    # ------------------------------------------------------------------

    def _render_bullet(self, slide, items: list, *,
                       top: float | None = None,
                       left: float | None = None,
                       width: float | None = None):
        """Render body items (strings as bullets, TableData as tables)."""
        t = self.theme
        top = top if top is not None else t.layout.body_top
        left = left if left is not None else self._margin.left
        width = width if width is not None else self._usable_width

        current_top = top
        for item in items:
            if isinstance(item, TableData):
                current_top = self._render_table(slide, item, current_top, left, width)
            else:
                current_top = self._render_text_bullet(
                    slide, str(item), current_top, left, width,
                )

    def _render_text_bullet(self, slide, text: str, top: float,
                            left: float, width: float) -> float:
        t = self.theme
        # Detect indent level
        level = 0
        clean = text
        if text.startswith("  ") or clean.strip().startswith("○"):
            level = 1
        clean = re.sub(r"^[\s●○\-]+", "", text).strip()
        if not clean:
            return top + 0.1  # small gap for empty lines

        indent = level * t.layout.bullet_indent
        actual_left = left + indent
        actual_width = width - indent

        # Bullet prefix styling
        bullet_char = "  " if level > 0 else ""
        if level == 0 and not text.strip().startswith(("🧪", "📌", "📊", "🚀", "💡", "🎯", "🖱", "☐")):
            bullet_char = "  "

        txbox = self._add_textbox(slide, actual_left, top, actual_width, 0.35)
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # Check for emphasis (arrows, key phrases)
        is_emphasis = "→" in clean or clean.startswith(("핵심", "정리"))

        run = p.add_run()
        run.text = clean
        run.font.name = t.fonts.body.name
        run.font.size = Pt(t.fonts.body.size)

        if is_emphasis:
            r, g, b = t.colors.hex_to_rgb("accent")
            run.font.color.rgb = RGBColor(r, g, b)
            run.font.bold = True
        else:
            r, g, b = t.colors.hex_to_rgb("secondary")
            run.font.color.rgb = RGBColor(r, g, b)

        return top + BULLET_HEIGHT

    def _render_table(self, slide, table: TableData, top: float,
                      left: float, width: float) -> float:
        t = self.theme
        rows = len(table.rows) + 1
        cols = len(table.headers)
        if cols == 0:
            return top

        tbl_shape = slide.shapes.add_table(
            rows, cols,
            Inches(left), Inches(top),
            Inches(width), Inches(TABLE_ROW_HEIGHT * rows),
        )
        tbl = tbl_shape.table

        # Style header row
        header_rgb = t.colors.hex_to_rgb("primary")
        for ci, header in enumerate(table.headers):
            cell = tbl.cell(0, ci)
            cell.text = header
            # Header background
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*header_rgb)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(t.fonts.body.size - 1)
                    run.font.name = t.fonts.body.name
                    run.font.color.rgb = RGBColor(255, 255, 255)

        # Data rows with alternating subtle background
        for ri, row in enumerate(table.rows, start=1):
            for ci, val in enumerate(row):
                if ci < cols:
                    cell = tbl.cell(ri, ci)
                    cell.text = val
                    # Alternating row colors
                    if ri % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(t.fonts.body.size - 1)
                            run.font.name = t.fonts.body.name

        return top + TABLE_ROW_HEIGHT * rows + TABLE_PADDING

    def _fill_background(self, slide, color_name: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        r, g, b = self.theme.colors.hex_to_rgb(color_name)
        fill.fore_color.rgb = RGBColor(r, g, b)

    # ------------------------------------------------------------------
    # Slide Type Renderers
    # ------------------------------------------------------------------

    def _render_title_slide(self, slide, model: Slide) -> None:
        """Memphis-style title: accent bar, large title, decorative elements."""
        t = self.theme

        # Side decoration strip
        self._add_side_decoration(slide, "accent")

        # Decorative dots
        self._add_corner_dot(slide, 1.5, 1.2, 0.25, "accent_yellow")
        self._add_corner_dot(slide, 10.5, 5.5, 0.2, "accent_purple")
        self._add_corner_dot(slide, 2.0, 5.8, 0.15, "accent")

        # Large accent bar
        self._add_accent_bar(slide, 2.5, width=1.2, color_name="accent")

        # Title
        txbox = self._add_textbox(
            slide, self._margin.left, 2.7, self._usable_width * 0.7, 1.2,
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = model.title
        run.font.name = t.fonts.title.name
        run.font.size = Pt(34)
        run.font.bold = True
        r, g, b = t.colors.hex_to_rgb("primary")
        run.font.color.rgb = RGBColor(r, g, b)

        # Subtitle from body[0]
        if model.body:
            sub_box = self._add_textbox(
                slide, self._margin.left, 4.1, self._usable_width * 0.6, 0.5,
            )
            sf = sub_box.text_frame
            sf.word_wrap = True
            sp = sf.paragraphs[0]
            sp.alignment = PP_ALIGN.LEFT
            srun = sp.add_run()
            srun.text = str(model.body[0])
            self._set_font(srun, t.fonts.subtitle, t.colors.hex_to_rgb("secondary"))

        # Bottom accent line
        self._add_accent_bar(slide, self._slide_h - 0.5, width=self._usable_width, color_name="accent_yellow")

    def _render_content_slide(self, slide, model: Slide) -> None:
        """Generic content: title + body bullets + tables with Memphis accents."""
        # Side strip decoration (alternating colors based on slide number)
        colors = ["accent", "accent_yellow", "accent_purple"]
        self._add_side_decoration(slide, colors[model.number % 3])

        self._add_title(slide, model.title, slide_model=model)
        self._add_page_number(slide, model)

        if model.body:
            self._render_bullet(slide, model.body)

    def _render_summary_slide(self, slide, model: Slide) -> None:
        """Corporate blue background + white text, Memphis dots."""
        self._fill_background(slide, "primary")
        t = self.theme

        # Decorative dots on dark background
        self._add_corner_dot(slide, 0.8, 0.6, 0.2, "accent")
        self._add_corner_dot(slide, 11.5, 6.2, 0.25, "accent_yellow")

        # Accent bar
        bar = slide.shapes.add_shape(
            1, Inches(self._margin.left), Inches(2.2),
            Inches(0.8), Inches(0.06),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*t.colors.hex_to_rgb("accent"))
        bar.line.fill.background()

        # Title
        txbox = self._add_textbox(
            slide, self._margin.left, 2.5, self._usable_width, 1.0,
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = model.title
        run.font.name = t.fonts.title.name
        run.font.size = Pt(t.fonts.title.size + 2)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Body items
        if model.body:
            body_top = 3.8
            for item in model.body:
                text = str(item)
                txb = self._add_textbox(
                    slide, self._margin.left + 0.1, body_top, self._usable_width - 0.2, 0.35,
                )
                btf = txb.text_frame
                btf.word_wrap = True
                bp = btf.paragraphs[0]
                br = bp.add_run()
                br.text = text
                br.font.name = t.fonts.body.name
                br.font.size = Pt(t.fonts.body.size)
                br.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
                body_top += 0.38

    def _render_exercise_slide(self, slide, model: Slide) -> None:
        """Checklist style with purple accent."""
        self._add_side_decoration(slide, "accent_purple")
        self._add_title(slide, model.title, slide_model=model)
        self._add_page_number(slide, model)

        t = self.theme
        top = t.layout.body_top
        accent_rgb = t.colors.hex_to_rgb("accent_purple")

        for item in model.body:
            text = str(item)
            clean = re.sub(r"^[\s●○\-🧪]+", "", text).strip()
            if not clean:
                continue

            # Checkbox with purple accent
            txbox = self._add_textbox(
                slide, self._margin.left, top, self._usable_width, 0.35,
            )
            tf = txbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]

            # Checkbox character
            check_run = p.add_run()
            check_run.text = "\u2610  "
            check_run.font.name = t.fonts.body.name
            check_run.font.size = Pt(t.fonts.body.size)
            check_run.font.color.rgb = RGBColor(*accent_rgb)

            # Text
            text_run = p.add_run()
            text_run.text = clean
            text_run.font.name = t.fonts.body.name
            text_run.font.size = Pt(t.fonts.body.size)
            r, g, b = t.colors.hex_to_rgb("secondary")
            text_run.font.color.rgb = RGBColor(r, g, b)

            top += BULLET_HEIGHT

    def _render_example_slide(self, slide, model: Slide) -> None:
        """2-column with question/insight split and accent divider."""
        self._add_side_decoration(slide, "accent_yellow")
        self._add_title(slide, model.title, slide_model=model)
        self._add_page_number(slide, model)

        half = self._usable_width / 2 - 0.2
        left_items = []
        right_items = []
        current = left_items

        for item in model.body:
            text = str(item) if not isinstance(item, TableData) else item
            if isinstance(text, str):
                low = text.lower()
                if any(kw in low for kw in ("인사이트", "insight", "결과", "액션")):
                    current = right_items
                elif any(kw in low for kw in ("질문", "question", "상황")):
                    current = left_items
            current.append(text)

        if not right_items:
            mid = len(model.body) // 2
            left_items = list(model.body[:mid])
            right_items = list(model.body[mid:])

        # Column headers
        t = self.theme
        col_header_top = t.layout.body_top - 0.05

        # Left column header
        lh = self._add_textbox(slide, self._margin.left, col_header_top, half, 0.3)
        lhp = lh.text_frame.paragraphs[0]
        lhr = lhp.add_run()
        lhr.text = "Question"
        lhr.font.name = t.fonts.body.name
        lhr.font.size = Pt(t.fonts.body.size - 1)
        lhr.font.bold = True
        lhr.font.color.rgb = RGBColor(*t.colors.hex_to_rgb("accent"))

        # Right column header
        rh = self._add_textbox(slide, self._margin.left + half + 0.4, col_header_top, half, 0.3)
        rhp = rh.text_frame.paragraphs[0]
        rhr = rhp.add_run()
        rhr.text = "Insight"
        rhr.font.name = t.fonts.body.name
        rhr.font.size = Pt(t.fonts.body.size - 1)
        rhr.font.bold = True
        rhr.font.color.rgb = RGBColor(*t.colors.hex_to_rgb("accent_purple"))

        # Vertical divider line
        divider = slide.shapes.add_shape(
            1,
            Inches(self._margin.left + half + 0.15),
            Inches(t.layout.body_top + 0.2),
            Inches(0.03),
            Inches(self._slide_h - t.layout.body_top - 1.0),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        divider.line.fill.background()

        body_top = t.layout.body_top + 0.3
        self._render_bullet(slide, left_items, left=self._margin.left, width=half, top=body_top)
        self._render_bullet(slide, right_items, left=self._margin.left + half + 0.4, width=half, top=body_top)

    def _render_two_column_slide(self, slide, model: Slide) -> None:
        """CRM/Experiment 2-column with accent divider."""
        self._add_side_decoration(slide, "accent")
        self._add_title(slide, model.title, slide_model=model)
        self._add_page_number(slide, model)

        t = self.theme
        half = self._usable_width / 2 - 0.2

        # Split body
        left_items = []
        right_items = []
        current = left_items
        for item in model.body:
            text = str(item) if not isinstance(item, TableData) else ""
            if "실험" in text or "Experiment" in text:
                current = right_items
            elif "CRM" in text or "🚀" in text:
                current = left_items
            current.append(item)

        if not right_items:
            mid = len(model.body) // 2
            left_items = list(model.body[:mid])
            right_items = list(model.body[mid:])

        # Column headers
        col_top = t.layout.body_top - 0.05
        lh = self._add_textbox(slide, self._margin.left, col_top, half, 0.3)
        lhr = lh.text_frame.paragraphs[0].add_run()
        lhr.text = "CRM 실행"
        lhr.font.name = t.fonts.body.name
        lhr.font.size = Pt(t.fonts.body.size - 1)
        lhr.font.bold = True
        lhr.font.color.rgb = RGBColor(*t.colors.hex_to_rgb("accent"))

        rh = self._add_textbox(slide, self._margin.left + half + 0.4, col_top, half, 0.3)
        rhr = rh.text_frame.paragraphs[0].add_run()
        rhr.text = "실험 설계"
        rhr.font.name = t.fonts.body.name
        rhr.font.size = Pt(t.fonts.body.size - 1)
        rhr.font.bold = True
        rhr.font.color.rgb = RGBColor(*t.colors.hex_to_rgb("accent_purple"))

        # Divider
        divider = slide.shapes.add_shape(
            1,
            Inches(self._margin.left + half + 0.15),
            Inches(t.layout.body_top + 0.2),
            Inches(0.03),
            Inches(self._slide_h - t.layout.body_top - 1.0),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        divider.line.fill.background()

        body_top = t.layout.body_top + 0.3
        self._render_bullet(slide, left_items, left=self._margin.left, width=half, top=body_top)
        self._render_bullet(slide, right_items, left=self._margin.left + half + 0.4, width=half, top=body_top)

    def _render_image_placeholder(self, slide, model: Slide) -> None:
        """Gray box placeholder with Memphis styling."""
        self._add_side_decoration(slide, "accent_yellow")
        self._add_title(slide, model.title, slide_model=model)
        t = self.theme

        box_w = self._usable_width * 0.8
        box_h = self._slide_h * 0.55
        box_left = self._margin.left + (self._usable_width - box_w) / 2
        box_top = t.layout.body_top + 0.2

        shape = slide.shapes.add_shape(
            1,
            Inches(box_left), Inches(box_top),
            Inches(box_w), Inches(box_h),
        )
        fill = shape.fill
        fill.solid()
        r, g, b = t.colors.hex_to_rgb("placeholder_bg")
        fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()

        # Description
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        if model.body:
            run = p.add_run()
            run.text = str(model.body[0])
            self._set_font(run, t.fonts.body, t.colors.hex_to_rgb("secondary"))

        # Caption
        cap_box = self._add_textbox(
            slide, box_left, box_top + box_h + 0.1, box_w, 0.3,
        )
        cf = cap_box.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = "이미지 수동 삽입 필요"
        self._set_font(cr, t.fonts.note, t.colors.hex_to_rgb("secondary"))

    def _add_section_divider(self, pptx, layout, section_title: str) -> None:
        """Corporate blue background + white centered section title with Memphis."""
        slide = pptx.slides.add_slide(layout)
        self._fill_background(slide, "section_bg")
        t = self.theme

        # Decorative elements
        self._add_corner_dot(slide, 1.0, 1.0, 0.3, "accent")
        self._add_corner_dot(slide, 10.0, 5.5, 0.25, "accent_yellow")

        # Accent bar
        bar = slide.shapes.add_shape(
            1, Inches(self._margin.left), Inches(3.0),
            Inches(1.0), Inches(0.06),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*t.colors.hex_to_rgb("accent"))
        bar.line.fill.background()

        # Section title
        txbox = self._add_textbox(
            slide, self._margin.left, 3.3, self._usable_width, 1.0,
        )
        tf = txbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = section_title
        run.font.name = t.fonts.title.name
        run.font.size = Pt(32)
        run.font.bold = True
        r, g, b = t.colors.hex_to_rgb("section_text")
        run.font.color.rgb = RGBColor(r, g, b)

        # Bottom yellow bar
        bottom_bar = slide.shapes.add_shape(
            1, Inches(0), Inches(self._slide_h - 0.12),
            Inches(self._slide_w), Inches(0.12),
        )
        bottom_bar.fill.solid()
        bottom_bar.fill.fore_color.rgb = RGBColor(*t.colors.hex_to_rgb("accent_yellow"))
        bottom_bar.line.fill.background()
